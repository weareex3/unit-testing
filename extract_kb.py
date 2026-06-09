"""
Extract text from .pptx and .docx files in the implementation folder
and save as KB_XX_name.txt files, continuing from the existing KB_ series.
"""

import os
import re
from pathlib import Path

SOURCE = Path(r"G:\My Drive\implementation")

def next_kb_number(folder):
    existing = [f.name for f in folder.glob("KB_*.txt")]
    nums = []
    for name in existing:
        m = re.match(r"KB_(\d+)_", name)
        if m:
            nums.append(int(m.group(1)))
    return max(nums) + 1 if nums else 1

def clean_name(stem):
    s = re.sub(r"[^\w]", "_", stem)
    s = re.sub(r"_+", "_", s).strip("_")
    return s[:60]

def already_extracted(folder, stem):
    target = clean_name(stem)
    for f in folder.glob("KB_*.txt"):
        if target.lower() in f.name.lower():
            return True
    return False

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    lines = [f"Source: {path.name}", ""]
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
        if slide_lines:
            lines.append(f"--- Slide {i} ---")
            lines.extend(slide_lines)
            lines.append("")
    return "\n".join(lines)

def extract_docx(path):
    from docx import Document
    doc = Document(path)
    lines = [f"Source: {path.name}", ""]
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            lines.append(text)
    for table in doc.tables:
        lines.append("")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)

def run():
    counter = next_kb_number(SOURCE)
    targets = (
        list(SOURCE.glob("*.pptx")) +
        list(SOURCE.glob("*.docx"))
    )

    skipped = []
    done = []
    errors = []

    for path in sorted(targets):
        if already_extracted(SOURCE, path.stem):
            skipped.append(path.name)
            continue

        kb_name = f"KB_{counter:02d}_{clean_name(path.stem)}.txt"
        out_path = SOURCE / kb_name

        try:
            if path.suffix == ".pptx":
                content = extract_pptx(path)
            else:
                content = extract_docx(path)

            out_path.write_text(content, encoding="utf-8")
            done.append(f"  KB_{counter:02d} <- {path.name}")
            counter += 1

        except Exception as e:
            errors.append(f"  SKIP {path.name}: {e}")

    print(f"\nExtracted {len(done)} files:")
    for d in done:
        print(d)

    if skipped:
        print(f"\nAlready had KB equivalent ({len(skipped)} skipped):")
        for s in skipped:
            print(f"  {s}")

    if errors:
        print(f"\nErrors ({len(errors)}):")
        for e in errors:
            print(e)

if __name__ == "__main__":
    run()
